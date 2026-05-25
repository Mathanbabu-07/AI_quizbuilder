import asyncio
from dataclasses import dataclass
from io import BytesIO

from fastapi import HTTPException, status

from app.services.pdf_extractor import clean_extracted_text


@dataclass
class PPTExtractor:
    max_slides: int = 120
    max_characters: int = 90_000

    async def extract_text(self, payload: bytes) -> str:
        return await asyncio.to_thread(self._extract_sync, payload)

    def _extract_sync(self, payload: bytes) -> str:
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise HTTPException(
                status_code=status.HTTP_503_SERVICE_UNAVAILABLE,
                detail="PPTX extraction dependency is missing. Install python-pptx on the backend.",
            ) from exc

        try:
            presentation = Presentation(BytesIO(payload))
        except Exception as exc:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="Could not read this PPTX. Upload a valid PowerPoint file.",
            ) from exc

        slides: list[str] = []
        for slide_index, slide in enumerate(presentation.slides):
            if slide_index >= self.max_slides:
                break

            lines: list[str] = []
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if not text:
                    continue
                lines.extend(part.strip() for part in text.splitlines() if part.strip())

            if lines:
                slides.append(f"Slide {slide_index + 1}\n" + "\n".join(lines))

            if sum(len(chunk) for chunk in slides) >= self.max_characters:
                break

        return clean_extracted_text("\n\n".join(slides), self.max_characters)
